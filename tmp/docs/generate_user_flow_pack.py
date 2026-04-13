#!/usr/bin/env python3
from __future__ import annotations

import argparse
import math
from pathlib import Path
from textwrap import dedent

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from flow_definitions import (
    ASSETS_DIR,
    BASE_URL,
    DECKS_DIR,
    MANUALS_DIR,
    ROLE_MAP,
    VERIFIED_DATE,
    WORKFLOWS,
    asset_relative_path,
    deck_filename,
    manual_filename,
)


ROOT = Path(__file__).resolve().parents[2]
MANUALS_ROOT = ROOT / MANUALS_DIR
ASSETS_ROOT = ROOT / ASSETS_DIR
DECKS_ROOT = ROOT / DECKS_DIR

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x0F, 0x76, 0x6E)
ACCENT_2 = RGBColor(0x1D, 0x4E, 0xD8)
WARM = RGBColor(0xF5, 0x9E, 0x0B)
SURFACE = RGBColor(0xF8, 0xFA, 0xFC)
PANEL = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x16, 0xA3, 0x4A)
ALERT = RGBColor(0xDC, 0x26, 0x26)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def ensure_dirs() -> None:
    MANUALS_ROOT.mkdir(parents=True, exist_ok=True)
    ASSETS_ROOT.mkdir(parents=True, exist_ok=True)
    DECKS_ROOT.mkdir(parents=True, exist_ok=True)
    for workflow in WORKFLOWS:
        (ASSETS_ROOT / workflow["slug"]).mkdir(parents=True, exist_ok=True)


def existing_image(workflow: dict, filename: str) -> Path | None:
    path = ASSETS_ROOT / workflow["slug"] / filename
    return path if path.exists() else None


def first_image_for_workflow(workflow: dict) -> Path | None:
    for step in workflow["steps"]:
        path = existing_image(workflow, step["screenshot"]["filename"])
        if path:
            return path
    return None


def wrap_bullets(items: list[str]) -> str:
    return "\n".join(f"- {item}" for item in items)


def role_map_mermaid() -> str:
    lines = [
        "```mermaid",
        "flowchart LR",
        '    Doctor["Doctor"] --> Support["Role-based support centers"]',
        '    Clinic["Clinic Staff"] --> Support',
        '    Brand["Brand Manager"] --> Support',
        '    Rep["Field Rep"] --> Support',
        '    Patient["Patient"] --> Support',
        '    Support --> PM["Project Manager dashboard"]',
        '    PM --> Ticketing["Ticket queue"]',
        '    Ticketing --> Owner["Department owner / support lead"]',
        '    PM --> Performance["Campaign performance"]',
        "```",
    ]
    return "\n".join(lines)


def inventory_markdown() -> str:
    section_lines = []
    for section in sorted({workflow["section"] for workflow in WORKFLOWS}):
        section_lines.append(f"## {section}")
        for workflow in [w for w in WORKFLOWS if w["section"] == section]:
            section_lines.append(
                f"- `{workflow['order']:02d}` [{workflow['title']}]({manual_filename(workflow)})"
                f" | Primary user: {workflow['primary_user']}"
            )
        section_lines.append("")

    table = [
        "| ID | Workflow | Section | Primary User | Entry Point | Status |",
        "| --- | --- | --- | --- | --- | --- |",
    ]
    for workflow in WORKFLOWS:
        table.append(
            f"| {workflow['order']:02d} | {workflow['title']} | {workflow['section']} | {workflow['primary_user']} | "
            f"`{workflow['entry_point']}` | {workflow['status']} |"
        )

    role_lines = wrap_bullets([f"**{role}**: {description}" for role, description in ROLE_MAP.items()])
    table_lines = "\n".join(table)
    section_text = "\n".join(section_lines)

    return dedent(
        f"""\
        # Workflow Inventory

        This inventory is the top-level guide to the user-flow training pack generated from the live application and repository source of truth.

        ## Scope

        - Verification date: `{VERIFIED_DATE}`
        - Source of truth order:
          1. live application behavior
          2. application code and templates
          3. project handoff and extracted documents
          4. older manual text
        - Demo base URL used for capture: `{BASE_URL}`

        ## Role Map

        {role_map_mermaid()}

        ## Roles

        {role_lines}

        ## Workflow Inventory Table

        {table_lines}

        {section_text}
        """
    ).strip() + "\n"


def workflow_markdown(workflow: dict) -> str:
    lines: list[str] = [f"# {workflow['title']}", ""]
    lines.extend(["## Document Purpose", "", workflow["purpose"], ""])
    lines.extend(["## Primary User", "", workflow["primary_user"], ""])
    lines.extend(["## Entry Point", "", f"`{workflow['entry_point']}`", ""])
    lines.extend(["## Workflow Summary", "", wrap_bullets(workflow["summary"]), ""])

    if workflow["slug"] == "platform-overview-and-role-map":
        lines.extend(
            [
                "## Role Diagram",
                "",
                role_map_mermaid(),
                "",
            ]
        )

    lines.extend(["## Step-By-Step Instructions", ""])

    for index, step in enumerate(workflow["steps"], start=1):
        relative_asset = asset_relative_path(workflow, step["screenshot"]["filename"])
        lines.extend(
            [
                f"### Step {index}. {step['title']}",
                "",
                f"- What the user does: {step['user_does']}",
                f"- What the user sees: {step['user_sees']}",
                f"- Why the step matters: {step['why']}",
                f"- Expected result: {step['expected_result']}",
                f"- Common issues or trainer notes: {step['notes']}",
                "- Screenshot placeholder:",
                f"  - Suggested file path: `{relative_asset}`",
                f"  - Screenshot caption: {step['screenshot']['caption']}",
                f"  - What the screenshot should show: {step['screenshot']['show']}",
                "",
            ]
        )
        asset_path = ASSETS_ROOT / workflow["slug"] / step["screenshot"]["filename"]
        if asset_path.exists():
            lines.extend([f"![{step['screenshot']['caption']}]({relative_asset})", ""])

    lines.extend(["## Success Criteria", "", wrap_bullets(workflow["success_criteria"]), ""])
    lines.extend(["## Related Documents", "", wrap_bullets([f"`{item}`" for item in workflow["related_documents"]]), ""])
    lines.extend(["## Status", "", workflow["status"], ""])
    return "\n".join(lines).strip() + "\n"


def write_manuals() -> list[Path]:
    written: list[Path] = []
    inventory_path = MANUALS_ROOT / "00-workflow-inventory.md"
    inventory_path.write_text(inventory_markdown(), encoding="utf-8")
    written.append(inventory_path)

    for workflow in WORKFLOWS:
        path = MANUALS_ROOT / manual_filename(workflow)
        path.write_text(workflow_markdown(workflow), encoding="utf-8")
        written.append(path)
    return written


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as image:
        return image.size


def image_ratio(path: Path | None) -> float:
    if not path:
        return 1.0
    width, height = image_size(path)
    return width / height


def add_fill(shape, color: RGBColor, line_color: RGBColor | None = None, line_width: Pt = Pt(1)) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_width


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=20,
    color=INK,
    bold=False,
    font=FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_paragraph(paragraph, text: str, *, size=16, color=INK, bold=False, level=0, bullet=False, font=FONT_BODY):
    paragraph.text = ""
    paragraph.level = level
    paragraph.bullet = bullet
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_chrome(slide, title: str, eyebrow: str, subtitle: str | None = None) -> None:
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    add_fill(bg, SURFACE)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.52))
    add_fill(band, ACCENT)

    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.78), Inches(2.1), Inches(0.38))
    add_fill(chip, WHITE, PANEL)
    add_textbox(
        slide,
        Inches(0.72),
        Inches(0.84),
        Inches(1.8),
        Inches(0.22),
        eyebrow.upper(),
        size=10,
        color=ACCENT,
        bold=True,
        font=FONT_BODY,
    )

    add_textbox(slide, Inches(0.58), Inches(1.28), Inches(8.9), Inches(0.65), title, size=26, color=INK, bold=True, font=FONT_HEAD)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(1.92), Inches(8.7), Inches(0.7), subtitle, size=12, color=SLATE)

    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.06))
    add_fill(footer, PANEL)

    add_textbox(
        slide,
        Inches(10.9),
        Inches(0.78),
        Inches(1.9),
        Inches(0.28),
        f"Verified {VERIFIED_DATE}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_card(slide, left, top, width, height, *, fill=WHITE, line=PANEL):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    add_fill(card, fill, line)
    return card


def add_metric_card(slide, left, top, width, height, label: str, value: str, tone: RGBColor = ACCENT) -> None:
    add_card(slide, left, top, width, height, fill=WHITE, line=PANEL)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.16), width - Inches(0.3), Inches(0.24), label.upper(), size=9, color=tone, bold=True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.48), width - Inches(0.3), height - Inches(0.56), value, size=14, color=INK, bold=True)


def add_image_contain(slide, path: Path, left, top, width, height, *, frame=True) -> None:
    if frame:
        add_card(slide, left, top, width, height, fill=WHITE, line=PANEL)
    img_w, img_h = image_size(path)
    scale = min(width / img_w, height / img_h)
    render_w = img_w * scale
    render_h = img_h * scale
    render_left = left + (width - render_w) / 2
    render_top = top + (height - render_h) / 2
    slide.shapes.add_picture(str(path), render_left, render_top, width=render_w, height=render_h)


def add_bullet_box(slide, left, top, width, height, title: str, bullets: list[str], tone: RGBColor = ACCENT) -> None:
    add_card(slide, left, top, width, height, fill=WHITE, line=PANEL)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.15), width - Inches(0.36), Inches(0.26), title, size=12, color=tone, bold=True)
    box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.45), width - Inches(0.36), height - Inches(0.55))
    frame = box.text_frame
    frame.word_wrap = True
    for idx, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        set_paragraph(paragraph, item, size=12, color=INK, bullet=True)


def add_cover_slide(prs: Presentation, workflow: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    add_fill(bg, SURFACE)

    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(4.45), SLIDE_H)
    add_fill(left_panel, ACCENT)
    accent_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4.2), Inches(0.6), Inches(0.22), Inches(6.2))
    add_fill(accent_panel, WARM)
    title_length = len(workflow["title"])
    if title_length > 38:
        title_size = 18
        title_height = Inches(2.35)
        summary_top = Inches(3.75)
    elif title_length > 28:
        title_size = 21
        title_height = Inches(1.95)
        summary_top = Inches(3.3)
    else:
        title_size = 25
        title_height = Inches(1.55)
        summary_top = Inches(3.0)
    add_textbox(slide, Inches(0.6), Inches(0.85), Inches(3.1), Inches(0.35), workflow["section"].upper(), size=11, color=WHITE, bold=True)
    add_textbox(
        slide,
        Inches(0.6),
        Inches(1.42),
        Inches(3.15),
        title_height,
        workflow["title"],
        size=title_size,
        color=WHITE,
        bold=True,
        font=FONT_HEAD,
    )
    add_textbox(slide, Inches(0.6), summary_top, Inches(3.2), Inches(1.2), workflow["purpose"], size=12, color=WHITE)

    add_metric_card(slide, Inches(0.6), Inches(5.05), Inches(1.45), Inches(1.05), "Primary user", workflow["primary_user"])
    add_metric_card(slide, Inches(2.15), Inches(5.05), Inches(1.55), Inches(1.05), "Entry point", workflow["entry_point"].replace(BASE_URL, ""))

    hero = first_image_for_workflow(workflow)
    if hero:
        add_image_contain(slide, hero, Inches(4.95), Inches(0.7), Inches(7.65), Inches(5.6))
    else:
        add_card(slide, Inches(4.95), Inches(0.7), Inches(7.65), Inches(5.6), fill=WHITE, line=PANEL)
        add_textbox(slide, Inches(5.3), Inches(2.95), Inches(7.0), Inches(0.6), "Screenshot slot", size=24, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(5.0), Inches(6.55), Inches(7.5), Inches(0.5), f"Verified against the live application on {VERIFIED_DATE}.", size=11, color=MUTED)


def add_overview_slide(prs: Presentation, workflow: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, workflow["title"], "Workflow overview", "A training-ready summary of the workflow entry point, ownership, and outcomes.")

    add_metric_card(slide, Inches(0.58), Inches(2.6), Inches(2.7), Inches(1.05), "Primary user", workflow["primary_user"])
    add_metric_card(slide, Inches(3.47), Inches(2.6), Inches(2.9), Inches(1.05), "Entry point", workflow["entry_point"].replace(BASE_URL, ""))
    add_metric_card(slide, Inches(6.56), Inches(2.6), Inches(2.7), Inches(1.05), "Verified", VERIFIED_DATE)
    add_metric_card(slide, Inches(9.45), Inches(2.6), Inches(2.9), Inches(1.05), "Steps", str(len(workflow["steps"])))

    add_bullet_box(slide, Inches(0.58), Inches(3.95), Inches(6.2), Inches(2.38), "Workflow summary", workflow["summary"], ACCENT)
    add_bullet_box(slide, Inches(6.95), Inches(3.95), Inches(5.4), Inches(2.38), "Success criteria", workflow["success_criteria"], ACCENT_2)


def add_decision_slide(prs: Presentation, workflow: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Decision Guide", workflow["title"], "Use this slide to explain how the workflow branches from entry to completion.")
    prompts = [
        ("Start here", workflow["steps"][0]["title"], workflow["steps"][0]["why"]),
        (
            "Continue when resolved",
            workflow["steps"][max(1, min(1, len(workflow["steps"]) - 1))]["title"],
            workflow["success_criteria"][0],
        ),
        (
            "Escalate or hand off",
            workflow["steps"][-1]["title"],
            workflow["success_criteria"][-1],
        ),
    ]
    positions = [Inches(0.8), Inches(4.55), Inches(8.3)]
    colors = [ACCENT, ACCENT_2, WARM]
    for (label, headline, body), left, tone in zip(prompts, positions, colors):
        add_card(slide, left, Inches(2.3), Inches(3.3), Inches(3.25), fill=WHITE, line=PANEL)
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.2), Inches(2.52), Inches(1.28), Inches(0.34))
        add_fill(band, tone, tone)
        add_textbox(slide, left + Inches(0.32), Inches(2.58), Inches(1.0), Inches(0.2), label.upper(), size=9, color=WHITE, bold=True)
        add_textbox(slide, left + Inches(0.22), Inches(3.0), Inches(2.8), Inches(0.6), headline, size=16, color=INK, bold=True, font=FONT_HEAD)
        add_textbox(slide, left + Inches(0.22), Inches(3.78), Inches(2.86), Inches(1.35), body, size=12, color=SLATE)


def add_step_slide(prs: Presentation, workflow: dict, step: dict, index: int) -> None:
    image_path = existing_image(workflow, step["screenshot"]["filename"])
    if image_path and image_ratio(image_path) >= 1.55:
        add_wide_step_slide(prs, workflow, step, index, image_path)
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, f"Step {index}. {step['title']}", workflow["title"], step["screenshot"]["caption"])

    add_bullet_box(
        slide,
        Inches(0.58),
        Inches(2.3),
        Inches(5.15),
        Inches(1.15),
        "What the user does",
        [step["user_does"]],
        ACCENT,
    )
    add_bullet_box(
        slide,
        Inches(0.58),
        Inches(3.62),
        Inches(5.15),
        Inches(1.15),
        "What the user sees",
        [step["user_sees"]],
        ACCENT_2,
    )
    add_bullet_box(
        slide,
        Inches(0.58),
        Inches(4.94),
        Inches(5.15),
        Inches(1.22),
        "Why it matters",
        [step["why"], step["expected_result"]],
        WARM,
    )

    if image_path:
        add_image_contain(slide, image_path, Inches(5.98), Inches(2.3), Inches(6.78), Inches(4.15))
    else:
        add_card(slide, Inches(5.98), Inches(2.3), Inches(6.78), Inches(4.15), fill=WHITE, line=PANEL)
        add_textbox(slide, Inches(6.25), Inches(4.05), Inches(6.2), Inches(0.35), "Screenshot not available", size=18, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

    note_box = slide.shapes.add_textbox(Inches(5.98), Inches(6.58), Inches(6.75), Inches(0.32))
    frame = note_box.text_frame
    p = frame.paragraphs[0]
    set_paragraph(p, f"Trainer note: {step['notes']}", size=11, color=SLATE)


def add_wide_step_slide(prs: Presentation, workflow: dict, step: dict, index: int, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, f"Step {index}. {step['title']}", workflow["title"], step["screenshot"]["caption"])
    add_image_contain(slide, image_path, Inches(0.58), Inches(2.25), Inches(12.15), Inches(3.65))
    add_bullet_box(slide, Inches(0.58), Inches(6.05), Inches(4.0), Inches(0.78), "What the user does", [step["user_does"]], ACCENT)
    add_bullet_box(slide, Inches(4.74), Inches(6.05), Inches(3.95), Inches(0.78), "Expected result", [step["expected_result"]], ACCENT_2)
    add_bullet_box(slide, Inches(8.84), Inches(6.05), Inches(3.89), Inches(0.78), "Trainer note", [step["notes"]], WARM)


def add_tips_slide(prs: Presentation, workflow: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Trainer Tips and Common Issues", workflow["title"], "Use this slide to coach facilitators on pacing, known caveats, and product reality.")
    tips = workflow.get("trainer_tips", ["No additional trainer tips were captured for this workflow."])
    add_bullet_box(slide, Inches(0.58), Inches(2.3), Inches(6.0), Inches(3.9), "Trainer tips", tips, ACCENT)
    add_bullet_box(slide, Inches(6.78), Inches(2.3), Inches(5.95), Inches(3.9), "Status and caveats", [workflow["status"]], ACCENT_2)


def add_closing_slide(prs: Presentation, workflow: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Close-Out", workflow["title"], "End the walkthrough by confirming success and pointing learners to the next reference.")
    add_bullet_box(slide, Inches(0.58), Inches(2.45), Inches(6.0), Inches(2.75), "Success criteria", workflow["success_criteria"], SUCCESS)
    add_bullet_box(
        slide,
        Inches(6.78),
        Inches(2.45),
        Inches(5.95),
        Inches(2.75),
        "Related documents",
        [item for item in workflow["related_documents"]],
        ACCENT_2,
    )
    add_textbox(slide, Inches(0.58), Inches(5.72), Inches(12.0), Inches(0.5), f"Source of truth: live app, then code/templates, then project docs. Verified on {VERIFIED_DATE}.", size=12, color=SLATE)


def add_role_map_slide(prs: Presentation) -> None:
    workflow = WORKFLOWS[0]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, "Role Map", workflow["title"], "How the implemented roles hand work off across the live product.")

    roles = list(ROLE_MAP.items())
    start_y = Inches(2.15)
    for idx, (role, description) in enumerate(roles[:5]):
        top = start_y + Inches(idx * 0.82)
        add_card(slide, Inches(0.7), top, Inches(2.2), Inches(0.58), fill=WHITE, line=PANEL)
        add_textbox(slide, Inches(0.9), top + Inches(0.12), Inches(1.8), Inches(0.25), role, size=12, color=INK, bold=True)

    add_card(slide, Inches(3.45), Inches(2.25), Inches(2.65), Inches(0.92), fill=WHITE, line=PANEL)
    add_textbox(slide, Inches(3.7), Inches(2.48), Inches(2.1), Inches(0.26), "Role-based support", size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(3.7), Inches(2.8), Inches(2.0), Inches(0.34), "Landing pages, FAQs, widgets, assistant", size=10, color=SLATE)

    add_card(slide, Inches(6.35), Inches(2.25), Inches(2.75), Inches(0.92), fill=WHITE, line=PANEL)
    add_textbox(slide, Inches(6.6), Inches(2.48), Inches(2.2), Inches(0.26), "PM dashboard", size=14, color=ACCENT_2, bold=True)
    add_textbox(slide, Inches(6.6), Inches(2.8), Inches(2.1), Inches(0.34), "Triage, review, escalation, performance", size=10, color=SLATE)

    add_card(slide, Inches(9.35), Inches(2.25), Inches(3.2), Inches(0.92), fill=WHITE, line=PANEL)
    add_textbox(slide, Inches(9.6), Inches(2.48), Inches(2.6), Inches(0.26), "Ticket execution", size=14, color=WARM, bold=True)
    add_textbox(slide, Inches(9.6), Inches(2.8), Inches(2.6), Inches(0.34), "Department queue, notes, status, routing", size=10, color=SLATE)

    for x in [Inches(2.92), Inches(5.98), Inches(9.08)]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, Inches(2.45), Inches(0.35), Inches(0.45))
        add_fill(arrow, PANEL, PANEL)

    add_bullet_box(
        slide,
        Inches(3.45),
        Inches(4.1),
        Inches(9.1),
        Inches(2.05),
        "What to tell learners",
        [
            "External roles start in public support centers, not in separate operational portals.",
            "Project Managers own triage, unresolved support reviews, and campaign performance monitoring.",
            "Department owners execute the ticket lifecycle inside the scoped queue and ticket detail pages.",
        ],
        ACCENT,
    )


def add_template_reference_deck() -> Path:
    prs = new_presentation()
    demo_workflow = WORKFLOWS[1]
    add_cover_slide(prs, demo_workflow)
    add_overview_slide(prs, demo_workflow)
    add_decision_slide(prs, demo_workflow)
    add_step_slide(prs, demo_workflow, demo_workflow["steps"][0], 1)
    add_step_slide(prs, demo_workflow, demo_workflow["steps"][1], 2)
    add_tips_slide(prs, demo_workflow)
    add_closing_slide(prs, demo_workflow)
    output = DECKS_ROOT / "00-training-pack-template-reference.pptx"
    prs.save(output)
    return output


def workflow_deck(workflow: dict) -> Path:
    prs = new_presentation()
    add_cover_slide(prs, workflow)
    add_overview_slide(prs, workflow)
    if workflow["slug"] == "platform-overview-and-role-map":
        add_role_map_slide(prs)
    else:
        add_decision_slide(prs, workflow)
    for index, step in enumerate(workflow["steps"], start=1):
        add_step_slide(prs, workflow, step, index)
    add_tips_slide(prs, workflow)
    add_closing_slide(prs, workflow)
    output = DECKS_ROOT / deck_filename(workflow)
    prs.save(output)
    return output


def add_index_card(slide, left, top, width, height, workflow: dict) -> None:
    card = add_card(slide, left, top, width, height, fill=WHITE, line=PANEL)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.14), top + Inches(0.16), Inches(0.72), Inches(0.3))
    add_fill(band, ACCENT, ACCENT)
    add_textbox(slide, left + Inches(0.28), top + Inches(0.21), Inches(0.36), Inches(0.18), f"{workflow['order']:02d}", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.58), width - Inches(0.32), Inches(0.55), workflow["title"], size=14, color=INK, bold=True, font=FONT_HEAD)
    add_textbox(slide, left + Inches(0.18), top + Inches(1.18), width - Inches(0.32), Inches(0.4), workflow["primary_user"], size=10, color=ACCENT_2, bold=True)
    add_textbox(slide, left + Inches(0.18), top + Inches(1.56), width - Inches(0.32), Inches(0.56), workflow["purpose"], size=9, color=SLATE)
    card.click_action.hyperlink.address = deck_filename(workflow)


def write_index_deck() -> Path:
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    add_fill(bg, SURFACE)
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.65))
    add_fill(top, ACCENT)
    add_textbox(slide, Inches(0.72), Inches(1.1), Inches(6.0), Inches(0.42), "User Flow Training Pack", size=28, color=INK, bold=True, font=FONT_HEAD)
    add_textbox(slide, Inches(0.72), Inches(1.62), Inches(7.0), Inches(0.55), "Open this deck first. Every card links to a sibling workflow deck in the same shared folder.", size=12, color=SLATE)
    add_bullet_box(
        slide,
        Inches(0.72),
        Inches(2.45),
        Inches(5.15),
        Inches(2.2),
        "Pack contents",
        [
            "1 platform overview deck",
            "5 PM and operations decks",
            "5 self-service role decks",
            "1 widget integration deck",
            "1 reusable template reference deck",
        ],
        ACCENT,
    )
    add_bullet_box(
        slide,
        Inches(6.05),
        Inches(2.45),
        Inches(6.45),
        Inches(2.2),
        "Verified on",
        [
            f"{VERIFIED_DATE} against the live local demo environment",
            "Relative file hyperlinks target sibling `.pptx` files in this folder",
            "QA PDFs are generated separately for visual review",
        ],
        ACCENT_2,
    )
    template_card = add_card(slide, Inches(0.72), Inches(5.1), Inches(3.3), Inches(1.25), fill=WHITE, line=PANEL)
    add_textbox(slide, Inches(0.95), Inches(5.35), Inches(2.6), Inches(0.28), "Template reference deck", size=14, color=INK, bold=True, font=FONT_HEAD)
    add_textbox(slide, Inches(0.95), Inches(5.74), Inches(2.8), Inches(0.3), "Open the reusable slide-style reference deck.", size=10, color=SLATE)
    template_card.click_action.hyperlink.address = "00-training-pack-template-reference.pptx"

    grouped = {}
    for workflow in WORKFLOWS:
        grouped.setdefault(workflow["section"], []).append(workflow)

    for section, workflows in grouped.items():
        section_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_chrome(section_slide, section, "Index", "Open any workflow card below to jump to the corresponding deck file.")
        add_textbox(section_slide, Inches(0.58), Inches(2.0), Inches(6.5), Inches(0.4), f"{len(workflows)} linked deck(s)", size=12, color=MUTED)
        cols = 2
        card_w = Inches(5.95)
        card_h = Inches(1.55)
        start_left = Inches(0.7)
        start_top = Inches(2.45)
        gap_x = Inches(0.45)
        gap_y = Inches(0.38)
        for idx, workflow in enumerate(workflows):
            row = idx // cols
            col = idx % cols
            left = start_left + col * (card_w + gap_x)
            top = start_top + row * (card_h + gap_y)
            add_index_card(section_slide, left, top, card_w, card_h, workflow)

    output = DECKS_ROOT / "00-user-flow-training-index.pptx"
    prs.save(output)
    return output


def write_decks() -> list[Path]:
    outputs = [add_template_reference_deck(), write_index_deck()]
    for workflow in WORKFLOWS:
        outputs.append(workflow_deck(workflow))
    return outputs


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate workflow manuals and PowerPoint decks.")
    parser.add_argument("--manuals-only", action="store_true", help="Generate only the Markdown manuals.")
    parser.add_argument("--decks-only", action="store_true", help="Generate only the PowerPoint decks.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    ensure_dirs()
    if not args.decks_only:
        write_manuals()
    if not args.manuals_only:
        write_decks()


if __name__ == "__main__":
    main()
